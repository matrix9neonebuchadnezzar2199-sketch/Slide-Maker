"""renderer の描画回帰テスト。"""

from __future__ import annotations

from pptx import Presentation

import renderer


def test_content_markdown_table_renders_as_powerpoint_table() -> None:
    """content に漏れた Markdown 表も PowerPoint table として描画される。"""
    prs = Presentation()
    renderer.render_content(
        prs,
        {
            "type": "content",
            "title": "店舗数表",
            "points": [
                "|国|5月|",
                "|---|---|",
                "|日本|10|",
                "|米国|20|",
            ],
        },
    )

    slide = prs.slides[0]
    assert any(shape.has_table for shape in slide.shapes)
    assert not any("|---|" in getattr(shape, "text", "") for shape in slide.shapes)


def test_content_bullet_marker_stays_with_text() -> None:
    """箇条点だけが独立したテキストボックスにならない。"""
    prs = Presentation()
    renderer.render_content(
        prs,
        {
            "type": "content",
            "title": "トピックス",
            "points": ["◎5月度トピックス 当月は順調に推移"],
        },
    )

    texts = [getattr(shape, "text", "") for shape in prs.slides[0].shapes]
    assert "•" not in texts
    assert any("• ◎5月度トピックス" in text for text in texts)
