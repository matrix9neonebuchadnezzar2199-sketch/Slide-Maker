"""PPTX 描画エンジン — RENDERERS ディスパッチ方式."""

from __future__ import annotations

import logging
import re
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

import schema
from table_parser import is_markdown_table_row, parse_markdown_table_block
import utils

logger = logging.getLogger(__name__)

_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_EMPHASIS_RE = re.compile(r"\[\[(.+?)\]\]")


def _fill_solid(shape, color: RGBColor) -> None:
    """図形を単色塗りにする。"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    font_size: Pt = schema.SIZE_BODY,
    bold: bool = False,
    color: RGBColor = schema.TEXT_MAIN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    """単純テキストボックスを追加する。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    schema.set_jp_font(run, size=font_size, color=color, bold=bold)
    return box


def _add_rich_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    font_size: Pt = schema.SIZE_BODY,
    color: RGBColor = schema.TEXT_MAIN,
) -> Any:
    """**太字** と [[強調]] を解釈するテキストボックス。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    for seg_text, seg_bold, seg_emphasis in _parse_inline_markup(text):
        if not seg_text:
            continue
        run = p.add_run()
        run.text = seg_text
        seg_color = schema.PRIMARY if seg_emphasis else color
        schema.set_jp_font(
            run,
            size=font_size,
            color=seg_color,
            bold=seg_bold or seg_emphasis,
        )

    return box


def _parse_inline_markup(text: str) -> list[tuple[str, bool, bool]]:
    """インライン記法をセグメントに分割する。"""
    segments: list[tuple[str, bool, bool]] = []
    pos = 0
    while pos < len(text):
        bold_m = _BOLD_RE.search(text, pos)
        emph_m = _EMPHASIS_RE.search(text, pos)
        candidates = [(m, "bold") for m in [bold_m] if m] + [(m, "emph") for m in [emph_m] if m]
        if not candidates:
            segments.append((text[pos:], False, False))
            break
        candidates.sort(key=lambda x: x[0].start())
        match, kind = candidates[0]
        if match.start() > pos:
            segments.append((text[pos:match.start()], False, False))
        inner = match.group(1)
        if kind == "bold":
            segments.append((inner, True, False))
        else:
            segments.append((inner, True, True))
        pos = match.end()
    return segments


def _blank_layout(prs: Presentation):
    """白紙に近いレイアウト（タイトルのみ）を返す。"""
    return prs.slide_layouts[6]


def _add_slide(prs: Presentation):
    """新規スライドを追加する。"""
    return prs.slides.add_slide(_blank_layout(prs))


def _draw_header(slide, s: dict[str, Any]) -> None:
    """共通タイトル帯（左アクセントバー + タイトル + subhead）を描画する。"""
    # 左縦アクセントバー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        schema.MARGIN_X - schema.ACCENT_BAR_W,
        schema.TITLE_Y,
        schema.ACCENT_BAR_W,
        schema.TITLE_H,
    )
    _fill_solid(bar, schema.ACCENT)

    title = s.get("title", "")
    if title:
        _add_textbox(
            slide,
            schema.MARGIN_X,
            schema.TITLE_Y,
            schema.CONTENT_W,
            schema.TITLE_H,
            title,
            font_size=schema.SIZE_TITLE,
            bold=True,
            color=schema.PRIMARY,
        )

    subhead = s.get("subhead")
    if subhead:
        _add_textbox(
            slide,
            schema.MARGIN_X,
            schema.SUBHEAD_Y,
            schema.CONTENT_W,
            Emu(400000),
            subhead,
            font_size=schema.SIZE_SUBHEAD,
            color=schema.TEXT_SUB,
        )


_draw_title_header = _draw_header  # 後方互換


def render_title(prs: Presentation, s: dict[str, Any]) -> None:
    """表紙スライドを描画する。"""
    slide = _add_slide(prs)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0),
        Emu(0),
        schema.SLIDE_W,
        schema.COVER_BAND_H,
    )
    _fill_solid(band, schema.PRIMARY)

    title = s.get("title", "")
    _add_textbox(
        slide,
        schema.MARGIN_X,
        Emu(2000000),
        schema.CONTENT_W,
        Emu(1200000),
        title,
        font_size=schema.SIZE_TITLE_COVER,
        bold=True,
        color=schema.TEXT_ON_FILL,
    )

    date_str = s.get("date", "")
    if date_str:
        _add_textbox(
            slide,
            schema.MARGIN_X,
            Emu(5200000),
            schema.CONTENT_W,
            Emu(400000),
            date_str,
            font_size=schema.SIZE_CAPTION,
            color=schema.TEXT_SUB,
        )


def render_section(prs: Presentation, s: dict[str, Any]) -> None:
    """章扉スライドを描画する。"""
    slide = _add_slide(prs)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), schema.SLIDE_W, schema.SLIDE_H,
    )
    _fill_solid(bg, schema.PRIMARY)

    section_no = s.get("sectionNo")
    if section_no is not None:
        _add_textbox(
            slide,
            Emu(800000),
            Emu(800000),
            Emu(10000000),
            Emu(5000000),
            str(section_no),
            font_size=schema.SIZE_SECTION_NO,
            bold=True,
            color=schema.SECTION_NO_FILL,
            align=PP_ALIGN.LEFT,
        )

    title = s.get("title", "")
    _add_textbox(
        slide,
        schema.MARGIN_X,
        schema.SECTION_TITLE_Y,
        schema.CONTENT_W,
        Emu(1500000),
        title,
        font_size=schema.SIZE_TITLE_SECTION,
        bold=True,
        color=schema.TEXT_ON_FILL,
    )


def render_content(prs: Presentation, s: dict[str, Any]) -> None:
    """箇条書き本文スライドを描画する。"""
    points = [str(p) for p in (s.get("points") or [])]
    if _render_markdown_table_content_if_needed(prs, s, points):
        return

    slide = _add_slide(prs)
    _draw_header(slide, s)

    gap = schema.CARD_GAP
    if s.get("twoColumn") and s.get("columns"):
        cols = s["columns"]
        col_w = (schema.CONTENT_W - gap) // 2
        for ci, col_items in enumerate(cols[:2]):
            left = schema.MARGIN_X + ci * (col_w + gap)
            _draw_bullet_list(slide, left, schema.BODY_Y, col_w, col_items or [])
    else:
        _draw_bullet_list(slide, schema.MARGIN_X, schema.BODY_Y, schema.CONTENT_W, points)


def _render_markdown_table_content_if_needed(
    prs: Presentation,
    s: dict[str, Any],
    points: list[str],
) -> bool:
    """content に漏れた Markdown 表を、本物の PowerPoint 表として描画する。"""
    table_lines = [line for line in points if is_markdown_table_row(line)]
    if len(table_lines) < 2:
        return False

    # 表以外の文章が混じる場合は content として扱い、誤変換を避ける。
    if len(table_lines) != len([line for line in points if line.strip()]):
        return False

    parsed = parse_markdown_table_block(table_lines)
    if not parsed:
        return False

    table_slide = dict(parsed)
    table_slide["title"] = s.get("title") or "データ一覧"
    render_table(prs, table_slide)
    return True


def _draw_bullet_list(slide, left: Emu, top: Emu, width: Emu, items: list[str]) -> None:
    """箇条書きリストを描画する。"""
    for i, item in enumerate(items):
        y = top + i * schema.BULLET_LINE_H
        _add_bullet_rich_textbox(slide, left, y, width, schema.BULLET_LINE_H, str(item))


def _add_bullet_rich_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
) -> Any:
    """箇条点と本文を同一段落内に描画する。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    bullet_run = p.add_run()
    bullet_run.text = "• "
    schema.set_jp_font(bullet_run, size=schema.SIZE_BODY, color=schema.ACCENT, bold=True)

    for seg_text, seg_bold, seg_emphasis in _parse_inline_markup(text):
        if not seg_text:
            continue
        run = p.add_run()
        run.text = seg_text
        seg_color = schema.PRIMARY if seg_emphasis else schema.TEXT_MAIN
        schema.set_jp_font(
            run,
            size=schema.SIZE_BODY,
            color=seg_color,
            bold=seg_bold or seg_emphasis,
        )

    return box



def render_agenda(prs: Presentation, s: dict[str, Any]) -> None:
    """アジェンダスライドを描画する。"""
    slide = _add_slide(prs)
    _draw_header(slide, s)

    items = s.get("items") or []
    for i, item in enumerate(items):
        y = schema.BODY_Y + i * schema.AGENDA_LINE_H
        num_text = f"{i + 1:02d}"
        _add_textbox(
            slide, schema.MARGIN_X, y, Emu(600000), schema.AGENDA_LINE_H,
            num_text, font_size=schema.SIZE_BODY, bold=True, color=schema.ACCENT,
        )
        _add_textbox(
            slide, schema.MARGIN_X + Emu(700000), y,
            schema.CONTENT_W - Emu(700000), schema.AGENDA_LINE_H,
            item, font_size=schema.SIZE_BODY,
        )


def render_closing(prs: Presentation, s: dict[str, Any]) -> None:
    """結びスライドを描画する。"""
    slide = _add_slide(prs)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), schema.SLIDE_W, schema.SLIDE_H,
    )
    _fill_solid(bg, schema.PRIMARY)

    _add_textbox(
        slide,
        schema.MARGIN_X,
        schema.CLOSING_TEXT_Y,
        schema.CONTENT_W,
        Emu(1200000),
        "ご清聴ありがとうございました",
        font_size=schema.SIZE_TITLE_SECTION,
        bold=True,
        color=schema.TEXT_ON_FILL,
        align=PP_ALIGN.CENTER,
    )


def _fill_solid_border(shape, fill: RGBColor, border: RGBColor = schema.BORDER) -> None:
    """図形を単色塗り＋枠線にする。"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)


def _add_rounded_card(slide, left, top, width, height):
    """角丸カード図形を追加する。"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    _fill_solid_border(card, schema.BG_LIGHT)
    return card


def render_kpi(prs: Presentation, s: dict[str, Any]) -> None:
    """KPI カードスライドを描画する。"""
    slide = _add_slide(prs)
    _draw_header(slide, s)

    items = (s.get("items") or [])[: schema.MAX_COUNT["kpi"]]
    if not items:
        return

    cols = s.get("columns")
    if cols is None:
        cols = len(items)
    cols = max(1, min(int(cols), len(items), schema.MAX_COUNT["kpi"]))

    gap = schema.CARD_GAP
    card_w = (schema.CONTENT_W - gap * (cols - 1)) // cols
    card_h = schema.KPI_CARD_H
    area_h = schema.BODY_H
    card_y = schema.BODY_Y + max(0, (area_h - card_h) // 2)

    for i, item in enumerate(items[:cols]):
        x = schema.MARGIN_X + i * (card_w + gap)
        card = _add_rounded_card(slide, x, card_y, card_w, card_h)

        status = item.get("status", "neutral")
        status_color = schema.STATUS_COLORS.get(status, schema.STATUS_NEUTRAL)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, card_y, card_w, schema.KPI_STATUS_BAR_H,
        )
        _fill_solid(bar, status_color)

        pad = Emu(150000)
        inner_y = card_y + schema.KPI_STATUS_BAR_H + pad
        label = str(item.get("label", ""))[:14]
        value = str(item.get("value", ""))
        change = str(item.get("change", ""))[:10]

        _add_textbox(
            slide, x + pad, inner_y, card_w - pad * 2, Emu(400000),
            label, font_size=schema.SIZE_KPI_LABEL, color=schema.TEXT_SUB,
        )
        _add_textbox(
            slide, x + pad, inner_y + Emu(450000), card_w - pad * 2, Emu(800000),
            value, font_size=schema.SIZE_KPI_VALUE, bold=True, color=schema.PRIMARY,
        )
        _add_textbox(
            slide, x + pad, inner_y + Emu(1300000), card_w - pad * 2, Emu(350000),
            change, font_size=schema.SIZE_CAPTION, color=status_color,
        )


def render_barCompare(prs: Presentation, s: dict[str, Any]) -> None:
    """棒グラフ比較スライドを描画する。"""
    slide = _add_slide(prs)
    _draw_header(slide, s)

    stats = (s.get("stats") or [])[: schema.BAR_COMPARE_MAX_ROWS]
    if not stats:
        return

    show_trends = bool(s.get("showTrends", False))
    label_w = int(schema.CONTENT_W * 22 // 100)
    bar_x = schema.MARGIN_X + label_w + Emu(100000)
    bar_w = schema.CONTENT_W - label_w - Emu(200000)
    row_h = schema.BODY_H // max(len(stats), 1)

    all_vals = []
    for stat in stats:
        all_vals.append(utils.parse_number(str(stat.get("leftValue", ""))))
        all_vals.append(utils.parse_number(str(stat.get("rightValue", ""))))
    max_val = max(all_vals) if all_vals else 1.0
    if max_val <= 0:
        max_val = 1.0

    # 凡例（右上）
    legend_y = schema.BODY_Y - Emu(200000)
    _add_textbox(
        slide, schema.MARGIN_X + schema.CONTENT_W - Emu(2800000), legend_y,
        Emu(1200000), Emu(200000), "A", font_size=schema.SIZE_CAPTION, color=schema.PRIMARY,
    )
    leg_a = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        schema.MARGIN_X + schema.CONTENT_W - Emu(3000000), legend_y + Emu(50000),
        Emu(150000), Emu(100000),
    )
    _fill_solid(leg_a, schema.PRIMARY)
    _add_textbox(
        slide, schema.MARGIN_X + schema.CONTENT_W - Emu(1400000), legend_y,
        Emu(1200000), Emu(200000), "B", font_size=schema.SIZE_CAPTION, color=schema.ACCENT,
    )
    leg_b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        schema.MARGIN_X + schema.CONTENT_W - Emu(1600000), legend_y + Emu(50000),
        Emu(150000), Emu(100000),
    )
    _fill_solid(leg_b, schema.ACCENT)

    thin = Emu(180000)
    gap_in_row = Emu(80000)

    for i, stat in enumerate(stats):
        row_y = schema.BODY_Y + i * row_h
        label = str(stat.get("label", ""))[:12]
        _add_textbox(
            slide, schema.MARGIN_X, row_y, label_w, row_h,
            label, font_size=schema.SIZE_BODY_SM, color=schema.TEXT_MAIN,
            align=PP_ALIGN.RIGHT,
        )

        left_val = utils.parse_number(str(stat.get("leftValue", "")))
        right_val = utils.parse_number(str(stat.get("rightValue", "")))
        left_len = int(bar_w * left_val / max_val) if max_val else 0
        right_len = int(bar_w * right_val / max_val) if max_val else 0

        bar1_y = row_y + Emu(80000)
        bar2_y = bar1_y + thin + gap_in_row

        if left_len > 0:
            b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bar1_y, left_len, thin)
            _fill_solid(b1, schema.PRIMARY)
            _add_textbox(
                slide, bar_x + left_len + Emu(50000), bar1_y, Emu(800000), thin,
                str(stat.get("leftValue", "")), font_size=schema.SIZE_CAPTION,
            )

        if right_len > 0:
            b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bar2_y, right_len, thin)
            _fill_solid(b2, schema.ACCENT)
            _add_textbox(
                slide, bar_x + right_len + Emu(50000), bar2_y, Emu(800000), thin,
                str(stat.get("rightValue", "")), font_size=schema.SIZE_CAPTION,
            )

        if show_trends:
            trend = stat.get("trend")
            tri_x = schema.MARGIN_X + schema.CONTENT_W - Emu(250000)
            tri_y = row_y + row_h // 2 - Emu(80000)
            if trend == "up":
                tri = slide.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE, tri_x, tri_y, Emu(160000), Emu(160000),
                )
                _fill_solid(tri, schema.STATUS_GOOD)
            elif trend == "down":
                tri = slide.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE, tri_x, tri_y, Emu(160000), Emu(160000),
                )
                tri.rotation = 180.0
                _fill_solid(tri, schema.STATUS_BAD)


def render_compare(prs: Presentation, s: dict[str, Any]) -> None:
    """対比2カラムスライドを描画する。"""
    slide = _add_slide(prs)
    _draw_header(slide, s)

    gap = schema.CARD_GAP
    col_w = (schema.CONTENT_W - gap) // 2
    left_items = (s.get("leftItems") or [])[: schema.COMPARE_MAX_ITEMS]
    right_items = (s.get("rightItems") or [])[: schema.COMPARE_MAX_ITEMS]

    columns = [
        (schema.MARGIN_X, s.get("leftTitle", ""), schema.PRIMARY, left_items),
        (schema.MARGIN_X + col_w + gap, s.get("rightTitle", ""), schema.PRIMARY_LT, right_items),
    ]

    header_h = schema.COMPARE_HEADER_H
    body_top = schema.BODY_Y + header_h + gap
    body_h = schema.BODY_H - header_h - gap

    for col_x, col_title, header_color, items in columns:
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, col_x, schema.BODY_Y, col_w, header_h,
        )
        _fill_solid(header, header_color)
        _add_textbox(
            slide, col_x, schema.BODY_Y, col_w, header_h,
            str(col_title), font_size=schema.SIZE_BODY, bold=True,
            color=schema.TEXT_ON_FILL, align=PP_ALIGN.CENTER,
        )

        body_card = _add_rounded_card(slide, col_x, body_top, col_w, body_h)
        line_h = Emu(500000)
        for i, item in enumerate(items):
            y = body_top + Emu(100000) + i * line_h
            oval = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                col_x + Emu(120000), y + Emu(80000),
                schema.COMPARE_BULLET_OVAL, schema.COMPARE_BULLET_OVAL,
            )
            _fill_solid(oval, schema.ACCENT)
            _add_textbox(
                slide, col_x + Emu(250000), y, col_w - Emu(350000), line_h,
                str(item), font_size=schema.SIZE_BODY_SM, color=schema.TEXT_MAIN,
            )


def render_table(prs: Presentation, s: dict[str, Any]) -> None:
    """表スライドを描画する。"""
    slide = _add_slide(prs)
    _draw_header(slide, s)

    headers = s.get("headers") or []
    rows = s.get("rows") or []
    if not headers:
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1
    row_h = schema.TABLE_ROW_H
    table_h = min(row_h * n_rows, schema.BODY_H)
    if row_h * n_rows > schema.BODY_H:
        row_h = schema.BODY_H // n_rows

    cell_font = (
        schema.SIZE_CAPTION if n_cols > schema.TABLE_DENSE_COLS else schema.SIZE_BODY_SM
    )

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        schema.MARGIN_X, schema.BODY_Y,
        schema.CONTENT_W, table_h,
    )
    table = table_shape.table

    col_w = schema.CONTENT_W // max(n_cols, 1)
    for ci in range(n_cols):
        table.columns[ci].width = col_w

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = schema.PRIMARY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                schema.set_jp_font(
                    run, size=cell_font, color=schema.TEXT_ON_FILL, bold=True,
                )

    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri + 1, ci)
            text = str(row[ci]) if ci < len(row) else ""
            cell.text = text
            bg = schema.BG_WHITE if ri % 2 == 0 else schema.BG_LIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    schema.set_jp_font(run, size=cell_font, color=schema.TEXT_MAIN)


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "content": render_content,
    "agenda": render_agenda,
    "closing": render_closing,
    "kpi": render_kpi,
    "barCompare": render_barCompare,
    "compare": render_compare,
    "table": render_table,
}


def build_pptx(slide_data: list[dict[str, Any]], out_path: str, template: str | None = None) -> None:
    """slideData から PPTX を生成する。

    Args:
        slide_data: 検証済みスライド配列。
        out_path: 出力 .pptx パス。
        template: 任意のテンプレートパス。
    """
    prs = Presentation(template) if template else Presentation()
    prs.slide_width = schema.SLIDE_W
    prs.slide_height = schema.SLIDE_H

    for s in slide_data:
        slide_type = s.get("type", "content")
        renderer_fn = RENDERERS.get(slide_type, render_content)
        try:
            renderer_fn(prs, s)
        except Exception as exc:
            logger.warning("skip slide type=%s err=%s", slide_type, exc)

    prs.save(out_path)
